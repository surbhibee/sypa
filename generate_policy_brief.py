#!/usr/bin/env python3
"""
Generate a Policy Brief PPTX for the IndiaAI SYPA.
Based on the Policy Brief Guidelines (AY 2025-2026) and example briefs.
All content sourced from the SYPA Final Draft.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──
NAVY       = RGBColor(0x0B, 0x1D, 0x3A)   # Dark navy for headers
DARK_BLUE  = RGBColor(0x14, 0x2D, 0x5E)   # Section headers
MID_BLUE   = RGBColor(0x1E, 0x56, 0xA0)   # Accent blue
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)  # Lighter accent
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)   # Light blue background
ORANGE     = RGBColor(0xE8, 0x6C, 0x00)   # India-themed accent / highlight
SAFFRON    = RGBColor(0xFF, 0x99, 0x33)    # Saffron for India theme
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
VERY_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

# Slide dimensions (widescreen 13.333 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=12,
                font_color=BLACK, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return txBox


def add_rich_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    """Add an empty textbox for rich content building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Remove default paragraph
    return txBox


def add_paragraph(text_frame, text, font_size=11, font_color=BLACK, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT, font_name="Calibri",
                  space_before=0, space_after=4, bullet=False, level=0):
    """Add a paragraph to an existing text frame."""
    # If the text frame already has content, add new paragraph
    if text_frame.paragraphs[0].text or len(text_frame.paragraphs) > 1:
        p = text_frame.add_paragraph()
    else:
        p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    if bullet:
        # Enable bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
    return p


def add_bullet_paragraph(text_frame, text, font_size=11, font_color=DARK_GRAY,
                         bold=False, font_name="Calibri", space_before=2, space_after=3,
                         bullet_char="\u2022"):
    """Add a bullet-point paragraph."""
    full_text = f"{bullet_char} {text}"
    return add_paragraph(text_frame, full_text, font_size=font_size, font_color=font_color,
                        bold=bold, font_name=font_name, space_before=space_before,
                        space_after=space_after)


def add_chart_placeholder(slide, left, top, width, height, title="[Chart Placeholder]",
                          description="Insert chart/graph here"):
    """Add a placeholder box for a chart/graph."""
    shape = add_shape(slide, left, top, width, height, fill_color=VERY_LIGHT, line_color=MED_GRAY, line_width=1.5)
    # Add dashed line look via text indication
    shape.line.dash_style = 2  # dash style

    # Add title text
    add_textbox(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                text=title, font_size=11, font_color=MID_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Add description
    add_textbox(slide, left + Inches(0.15), top + height/2 - Inches(0.3), width - Inches(0.3), Inches(0.6),
                text=description, font_size=10, font_color=MED_GRAY, italic=True,
                alignment=PP_ALIGN.CENTER)
    return shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE / OVERVIEW
# ═══════════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(blank_layout)

# Full background navy bar on top
add_shape(slide1, Inches(0), Inches(0), SLIDE_W, Inches(3.4), fill_color=NAVY)

# Saffron accent line
add_shape(slide1, Inches(0), Inches(3.4), SLIDE_W, Inches(0.06), fill_color=SAFFRON)

# Title
add_textbox(slide1, Inches(0.6), Inches(0.35), Inches(8.5), Inches(0.65),
            text="Unlocking India\u2019s AI Potential", font_size=32, font_color=WHITE,
            bold=True, font_name="Calibri")

# Subtitle
add_textbox(slide1, Inches(0.6), Inches(1.0), Inches(8.5), Inches(0.5),
            text="Evaluation of IndiaAI Mission\u2019s Common Compute Facility",
            font_size=18, font_color=SAFFRON, bold=False, font_name="Calibri")

# Author info
add_textbox(slide1, Inches(0.6), Inches(1.65), Inches(8), Inches(0.35),
            text="Prateek Pillai & Surbhi Bharadwaj  |  Harvard Kennedy School  |  February 2026",
            font_size=11, font_color=RGBColor(0xBB, 0xBB, 0xBB), font_name="Calibri")

# Key stat callout boxes in the navy header area (right side)
# Stat box 1
stat_box1 = add_rounded_shape(slide1, Inches(9.6), Inches(0.4), Inches(1.6), Inches(1.15),
                               fill_color=RGBColor(0x14, 0x2D, 0x5E), line_color=SAFFRON)
add_textbox(slide1, Inches(9.6), Inches(0.45), Inches(1.6), Inches(0.6),
            text="34,000+", font_size=26, font_color=SAFFRON, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(slide1, Inches(9.6), Inches(0.95), Inches(1.6), Inches(0.55),
            text="GPUs\nEmpanelled", font_size=10, font_color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Stat box 2
stat_box2 = add_rounded_shape(slide1, Inches(11.4), Inches(0.4), Inches(1.6), Inches(1.15),
                               fill_color=RGBColor(0x14, 0x2D, 0x5E), line_color=SAFFRON)
add_textbox(slide1, Inches(11.4), Inches(0.45), Inches(1.6), Inches(0.6),
            text="<50%", font_size=26, font_color=ORANGE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(slide1, Inches(11.4), Inches(0.95), Inches(1.6), Inches(0.55),
            text="Capacity\nUtilized", font_size=10, font_color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Stat box 3
stat_box3 = add_rounded_shape(slide1, Inches(9.6), Inches(1.75), Inches(1.6), Inches(1.15),
                               fill_color=RGBColor(0x14, 0x2D, 0x5E), line_color=SAFFRON)
add_textbox(slide1, Inches(9.6), Inches(1.8), Inches(1.6), Inches(0.6),
            text="$1.25B", font_size=26, font_color=SAFFRON, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(slide1, Inches(9.6), Inches(2.3), Inches(1.6), Inches(0.55),
            text="Mission\nBudget (5yr)", font_size=10, font_color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Stat box 4
stat_box4 = add_rounded_shape(slide1, Inches(11.4), Inches(1.75), Inches(1.6), Inches(1.15),
                               fill_color=RGBColor(0x14, 0x2D, 0x5E), line_color=SAFFRON)
add_textbox(slide1, Inches(11.4), Inches(1.8), Inches(1.6), Inches(0.6),
            text="<200", font_size=26, font_color=ORANGE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(slide1, Inches(11.4), Inches(2.3), Inches(1.6), Inches(0.55),
            text="Registered\nPortal Users", font_size=10, font_color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")

# ── Summary section (bottom half of slide 1) ──
# Left column: Summary
add_textbox(slide1, Inches(0.6), Inches(3.7), Inches(3.5), Inches(0.35),
            text="Summary", font_size=16, font_color=NAVY, bold=True, font_name="Calibri")

summary_box = add_rich_textbox(slide1, Inches(0.6), Inches(4.1), Inches(5.5), Inches(3.1))
tf = summary_box.text_frame
add_paragraph(tf, "India\u2019s IndiaAI Mission launched the Common Compute Facility (CCF) "
              "to expand access to subsidized GPUs for startups, researchers, and public "
              "institutions. Despite empanelling over 34,000 GPUs, less than half of this "
              "capacity has been utilized.", font_size=10.5, font_color=DARK_GRAY,
              space_after=6)
add_paragraph(tf, "Our analysis identifies two binding constraints: (1) low perceived "
              "returns to AI innovation among target users, and (2) administrative frictions "
              "in accessing approved compute allocations. Human capital gaps act as an "
              "ancillary constraint, while financial costs and infrastructure mismatch "
              "are non-binding.", font_size=10.5, font_color=DARK_GRAY, space_after=6)
add_paragraph(tf, "We recommend two priority reforms: advance government procurement "
              "commitments in high-priority domains and time-bound GPU provisioning "
              "guarantees to restore credibility and stimulate demand.",
              font_size=10.5, font_color=DARK_GRAY, bold=True, space_after=2)

# Right column: Key Findings + Policy Recommendations
# Key Findings box
findings_bg = add_rounded_shape(slide1, Inches(6.4), Inches(3.65), Inches(3.0), Inches(3.35),
                                 fill_color=LIGHT_BLUE, line_color=MID_BLUE)
add_textbox(slide1, Inches(6.6), Inches(3.75), Inches(2.6), Inches(0.35),
            text="Key Findings", font_size=14, font_color=NAVY, bold=True, font_name="Calibri")

findings_tb = add_rich_textbox(slide1, Inches(6.6), Inches(4.15), Inches(2.7), Inches(2.8))
tf_f = findings_tb.text_frame
add_bullet_paragraph(tf_f, "Low perceived returns to AI innovation are the primary constraint "
                     "to CCF uptake", font_size=9.5, font_color=DARK_GRAY, space_after=5)
add_bullet_paragraph(tf_f, "Administrative delays in GPU provisioning deter private-sector "
                     "users and reduce the applicant pool", font_size=9.5, font_color=DARK_GRAY, space_after=5)
add_bullet_paragraph(tf_f, "Human capital gaps amplify binding constraints but do not "
                     "independently explain low utilization", font_size=9.5, font_color=DARK_GRAY, space_after=5)
add_bullet_paragraph(tf_f, "Price competitiveness is not a binding factor\u2014CCF rates are "
                     "<50% of market, yet users prefer commercial providers", font_size=9.5,
                     font_color=DARK_GRAY, space_after=5)
add_bullet_paragraph(tf_f, "Only 2\u20136 serious language-model players exist in the country; "
                     "broader AI use cases need to be actively built", font_size=9.5,
                     font_color=DARK_GRAY, space_after=2)

# Policy Recommendations box
rec_bg = add_rounded_shape(slide1, Inches(9.65), Inches(3.65), Inches(3.35), Inches(3.35),
                            fill_color=RGBColor(0xEB, 0xF5, 0xFB), line_color=ACCENT_BLUE)
add_textbox(slide1, Inches(9.85), Inches(3.75), Inches(3.0), Inches(0.35),
            text="Policy Recommendations", font_size=14, font_color=NAVY, bold=True, font_name="Calibri")

rec_tb = add_rich_textbox(slide1, Inches(9.85), Inches(4.2), Inches(3.0), Inches(2.7))
tf_r = rec_tb.text_frame

add_paragraph(tf_r, "1  Advance Government Procurement", font_size=11, font_color=MID_BLUE,
              bold=True, space_after=2)
add_paragraph(tf_r, "Pre-commit to procuring AI solutions in "
              "priority public-sector domains (e.g., OCR of legacy records, multilingual "
              "interfaces, forecasting) to credibly signal downstream demand.",
              font_size=9.5, font_color=DARK_GRAY, space_after=8)

add_paragraph(tf_r, "2  Time-Bound GPU Provisioning", font_size=11, font_color=MID_BLUE,
              bold=True, space_after=2)
add_paragraph(tf_r, "Introduce clear, publicly stated service-level "
              "guarantees for compute delivery after approval, with automatic escalation "
              "mechanisms if timelines are breached.",
              font_size=9.5, font_color=DARK_GRAY, space_after=2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: CONTEXT & IMPORTANCE OF THE PROBLEM
# ═══════════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(blank_layout)

# Top bar
add_shape(slide2, Inches(0), Inches(0), SLIDE_W, Inches(0.85), fill_color=NAVY)
add_shape(slide2, Inches(0), Inches(0.85), SLIDE_W, Inches(0.04), fill_color=SAFFRON)

# Page header
add_textbox(slide2, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
            text="Context & Importance of the Problem", font_size=22, font_color=WHITE,
            bold=True, font_name="Calibri")
# Page number
add_textbox(slide2, Inches(12.0), Inches(0.2), Inches(1), Inches(0.4),
            text="2", font_size=14, font_color=SAFFRON, bold=True,
            alignment=PP_ALIGN.RIGHT, font_name="Calibri")

# ── Left column: The Policy Puzzle ──
add_textbox(slide2, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.35),
            text="The Policy Puzzle", font_size=16, font_color=NAVY, bold=True, font_name="Calibri")

# Highlighted question box
puzzle_bg = add_rounded_shape(slide2, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.65),
                               fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide2, Inches(0.65), Inches(1.6), Inches(5.7), Inches(0.55),
            text="How can IndiaAI increase utilization of CCF in support of the broader objective "
                 "of stimulating India\u2019s AI innovation ecosystem?",
            font_size=11, font_color=DARK_GRAY, bold=True, italic=True,
            alignment=PP_ALIGN.LEFT, font_name="Calibri")

# Context text
ctx_tb = add_rich_textbox(slide2, Inches(0.5), Inches(2.4), Inches(6.0), Inches(4.8))
tf_ctx = ctx_tb.text_frame

add_paragraph(tf_ctx, "IndiaAI Mission & Common Compute Facility", font_size=13,
              font_color=DARK_BLUE, bold=True, space_after=4)
add_bullet_paragraph(tf_ctx, "The IndiaAI Mission, housed under MeitY, was launched in March "
                     "2024 with INR 10,371 crore ($1.25B) over five years to catalyze AI innovation "
                     "in India.", font_size=10, font_color=DARK_GRAY, space_after=4)
add_bullet_paragraph(tf_ctx, "The Common Compute Facility (CCF) is the Mission\u2019s largest pillar, "
                     "accounting for 44% of the total allocation ($550M). It empanels private data-center "
                     "providers and offers subsidized GPU access to startups, researchers, and public "
                     "institutions.", font_size=10, font_color=DARK_GRAY, space_after=4)
add_bullet_paragraph(tf_ctx, "By November 2025, CCF had empanelled 34,000+ GPUs but allocated less "
                     "than 50% of monthly capacity \u2014 despite compute rates well below market "
                     "prices.", font_size=10, font_color=DARK_GRAY, space_after=8)

add_paragraph(tf_ctx, "Why This Matters", font_size=13,
              font_color=DARK_BLUE, bold=True, space_after=4)
add_bullet_paragraph(tf_ctx, "Inflection point for tech leadership: India has strong technical talent "
                     "and a vibrant startup ecosystem, but AI innovation is far more capital-intensive. "
                     "India has only 5 compute regions vs. 22 in China and 26 in the US.",
                     font_size=10, font_color=DARK_GRAY, space_after=4)
add_bullet_paragraph(tf_ctx, "Sovereign AI imperative: Chinese and US companies operate 90% of "
                     "public cloud AI facilities globally. India is pursuing indigenous language models "
                     "and government AI applications \u2014 requiring accessible domestic compute.",
                     font_size=10, font_color=DARK_GRAY, space_after=4)
add_bullet_paragraph(tf_ctx, "Distributional equity: Without public intervention, well-capitalized "
                     "organizations dominate. CCF aims to expand access for smaller startups, "
                     "researchers, and academics \u2014 but only if utilized.",
                     font_size=10, font_color=DARK_GRAY, space_after=2)

# ── Right column: CCF Structure + chart placeholders ──
add_textbox(slide2, Inches(6.85), Inches(1.15), Inches(6.0), Inches(0.35),
            text="How CCF Works", font_size=16, font_color=NAVY, bold=True, font_name="Calibri")

ccf_tb = add_rich_textbox(slide2, Inches(6.85), Inches(1.55), Inches(6.0), Inches(1.6))
tf_ccf = ccf_tb.text_frame
add_paragraph(tf_ccf, "The government empanels private data-center operators through competitive "
              "bidding. GPUs remain privately owned; public funds are disbursed only when actually "
              "used. End users apply via the IndiaAI portal and receive subsidies of up to 40% of "
              "the bid price.", font_size=10, font_color=DARK_GRAY, space_after=4)
add_paragraph(tf_ccf, "Target users: Foundation model developers, startups using open-source "
              "models, government departments deploying AI on sensitive datasets \u2014 actors who "
              "need their own compute, not just API access.",
              font_size=10, font_color=DARK_GRAY, space_after=2)

# Chart placeholder: CCF Flow diagram
add_chart_placeholder(slide2, Inches(6.85), Inches(3.3), Inches(5.8), Inches(1.7),
                      title="[Figure: CCF Structure & Flow of Compute]",
                      description="Insert diagram showing procurement \u2192 empanelment \u2192 allocation \u2192 end use flow")

# Chart placeholder: Utilization gap
add_chart_placeholder(slide2, Inches(6.85), Inches(5.2), Inches(5.8), Inches(1.7),
                      title="[Figure: CCF GPU Allocation vs. Capacity]",
                      description="Insert chart showing empanelled capacity vs. actual utilization over time")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: ANALYSIS & KEY FINDINGS
# ═══════════════════════════════════════════════════════════════════

slide3 = prs.slides.add_slide(blank_layout)

# Top bar
add_shape(slide3, Inches(0), Inches(0), SLIDE_W, Inches(0.85), fill_color=NAVY)
add_shape(slide3, Inches(0), Inches(0.85), SLIDE_W, Inches(0.04), fill_color=SAFFRON)

add_textbox(slide3, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
            text="Analysis & Key Findings", font_size=22, font_color=WHITE,
            bold=True, font_name="Calibri")
add_textbox(slide3, Inches(12.0), Inches(0.2), Inches(1), Inches(0.4),
            text="3", font_size=14, font_color=SAFFRON, bold=True,
            alignment=PP_ALIGN.RIGHT, font_name="Calibri")

# Methodology note
add_textbox(slide3, Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.4),
            text="Based on 20+ stakeholder interviews across policymakers, data-center providers, end users, "
                 "and industry researchers, complemented by administrative data and comparative policy analysis.",
            font_size=10, font_color=MED_GRAY, italic=True, font_name="Calibri")

# ── Constraint classification: 4 boxes across ──
constraint_y = Inches(1.65)
box_h = Inches(2.6)
box_w = Inches(2.95)
gap = Inches(0.2)

# Box 1: Binding - Low perceived returns
b1 = add_rounded_shape(slide3, Inches(0.5), constraint_y, box_w, box_h,
                        fill_color=WHITE, line_color=RED_ACCENT)
add_shape(slide3, Inches(0.5), constraint_y, box_w, Inches(0.4), fill_color=RED_ACCENT)
add_textbox(slide3, Inches(0.55), constraint_y + Inches(0.02), box_w - Inches(0.1), Inches(0.35),
            text="BINDING CONSTRAINT", font_size=9, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

b1_tb = add_rich_textbox(slide3, Inches(0.65), constraint_y + Inches(0.5), Inches(2.65), Inches(2.0))
tf_b1 = b1_tb.text_frame
add_paragraph(tf_b1, "Low Perceived Returns to AI Innovation", font_size=11,
              font_color=DARK_BLUE, bold=True, space_after=5)
add_bullet_paragraph(tf_b1, "Target users do not view investing in AI through dedicated "
                     "compute as value-generating", font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b1, "Ecosystem bifurcation: only 2\u20136 serious language-model players; "
                     "most firms use AI via APIs", font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b1, "Uncertainty around monetization\u2014not cost\u2014is the key deterrent",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b1, "\u201CPeople still think CCF is only for language models\u201D \u2014 data-center provider",
                     font_size=9, font_color=MED_GRAY, space_after=2)

# Box 2: Binding - Administrative frictions
b2_x = Inches(0.5) + box_w + gap
b2 = add_rounded_shape(slide3, b2_x, constraint_y, box_w, box_h,
                        fill_color=WHITE, line_color=RED_ACCENT)
add_shape(slide3, b2_x, constraint_y, box_w, Inches(0.4), fill_color=RED_ACCENT)
add_textbox(slide3, b2_x + Inches(0.05), constraint_y + Inches(0.02), box_w - Inches(0.1), Inches(0.35),
            text="BINDING CONSTRAINT", font_size=9, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

b2_tb = add_rich_textbox(slide3, b2_x + Inches(0.15), constraint_y + Inches(0.5), Inches(2.65), Inches(2.0))
tf_b2 = b2_tb.text_frame
add_paragraph(tf_b2, "Administrative Frictions in Access", font_size=11,
              font_color=DARK_BLUE, bold=True, space_after=5)
add_bullet_paragraph(tf_b2, "Approved users wait extended periods before GPUs are provisioned",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b2, "Private-sector users (especially startups) face longer, more uncertain timelines "
                     "than public entities",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b2, "Delays cause self-selection out of system\u2014 fewer than 200 portal registrations",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b2, "\u201CI had completely forgotten about the CCF approval\u201D \u2014 startup founder",
                     font_size=9, font_color=MED_GRAY, space_after=2)

# Box 3: Ancillary - Human Capital
b3_x = b2_x + box_w + gap
b3 = add_rounded_shape(slide3, b3_x, constraint_y, box_w, box_h,
                        fill_color=WHITE, line_color=ORANGE)
add_shape(slide3, b3_x, constraint_y, box_w, Inches(0.4), fill_color=ORANGE)
add_textbox(slide3, b3_x + Inches(0.05), constraint_y + Inches(0.02), box_w - Inches(0.1), Inches(0.35),
            text="ANCILLARY CONSTRAINT", font_size=9, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

b3_tb = add_rich_textbox(slide3, b3_x + Inches(0.15), constraint_y + Inches(0.5), Inches(2.65), Inches(2.0))
tf_b3 = b3_tb.text_frame
add_paragraph(tf_b3, "Human Capital Gaps", font_size=11,
              font_color=DARK_BLUE, bold=True, space_after=5)
add_bullet_paragraph(tf_b3, "Limited pool of professionals with GPU-intensive model experience",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b3, "Amplifies binding constraints: orgs lacking expertise are less likely to "
                     "pursue compute-intensive work",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b3, "Low public R&D spending constrains advanced research training",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b3, "May become more salient as utilization grows",
                     font_size=9, font_color=DARK_GRAY, space_after=2)

# Box 4: Non-Binding
b4_x = b3_x + box_w + gap
b4 = add_rounded_shape(slide3, b4_x, constraint_y, box_w, box_h,
                        fill_color=WHITE, line_color=GREEN)
add_shape(slide3, b4_x, constraint_y, box_w, Inches(0.4), fill_color=GREEN)
add_textbox(slide3, b4_x + Inches(0.05), constraint_y + Inches(0.02), box_w - Inches(0.1), Inches(0.35),
            text="NON-BINDING", font_size=9, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

b4_tb = add_rich_textbox(slide3, b4_x + Inches(0.15), constraint_y + Inches(0.5), Inches(2.65), Inches(2.0))
tf_b4 = b4_tb.text_frame
add_paragraph(tf_b4, "Price & Infrastructure", font_size=11,
              font_color=DARK_BLUE, bold=True, space_after=5)
add_bullet_paragraph(tf_b4, "CCF costs <50% of market rates \u2014 yet users prefer commercial "
                     "providers for non-price reasons",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b4, "Further price reductions unlikely to drive uptake",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b4, "CCF provides latest GPUs (incl. H100s); infrastructure mismatch not "
                     "an issue at current utilization",
                     font_size=9, font_color=DARK_GRAY, space_after=3)
add_bullet_paragraph(tf_b4, "GPUs are identical to those from private providers",
                     font_size=9, font_color=DARK_GRAY, space_after=2)

# ── Chart placeholders at bottom ──
add_chart_placeholder(slide3, Inches(0.5), Inches(4.55), Inches(3.9), Inches(2.55),
                      title="[Figure: Fishbone / Constraint Decomposition]",
                      description="Insert fishbone diagram showing demand- and supply-side\nfactors driving underutilization")

add_chart_placeholder(slide3, Inches(4.65), Inches(4.55), Inches(3.9), Inches(2.55),
                      title="[Figure: Constraint Evaluation Matrix]",
                      description="Insert matrix rating each constraint as\nbinding / ancillary / non-binding")

add_chart_placeholder(slide3, Inches(8.8), Inches(4.55), Inches(4.2), Inches(2.55),
                      title="[Figure: End User Allocation & Subsidy Distribution]",
                      description="Insert chart showing user categories, allocation sizes,\nand subsidy distribution")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: POLICY RECOMMENDATIONS & IMPLEMENTATION
# ═══════════════════════════════════════════════════════════════════

slide4 = prs.slides.add_slide(blank_layout)

# Top bar
add_shape(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.85), fill_color=NAVY)
add_shape(slide4, Inches(0), Inches(0.85), SLIDE_W, Inches(0.04), fill_color=SAFFRON)

add_textbox(slide4, Inches(0.5), Inches(0.15), Inches(10), Inches(0.55),
            text="Policy Recommendations & Implementation", font_size=22, font_color=WHITE,
            bold=True, font_name="Calibri")
add_textbox(slide4, Inches(12.0), Inches(0.2), Inches(1), Inches(0.4),
            text="4", font_size=14, font_color=SAFFRON, bold=True,
            alignment=PP_ALIGN.RIGHT, font_name="Calibri")

# ── Recommendation 1 (left half) ──
rec1_x = Inches(0.5)
rec1_w = Inches(6.1)

# Header box
r1_header = add_rounded_shape(slide4, rec1_x, Inches(1.15), rec1_w, Inches(0.55),
                               fill_color=MID_BLUE)
add_textbox(slide4, rec1_x + Inches(0.15), Inches(1.2), rec1_w - Inches(0.3), Inches(0.45),
            text="1  Advance Government Procurement in Priority Domains",
            font_size=14, font_color=WHITE, bold=True, font_name="Calibri")

# Addresses label
add_textbox(slide4, rec1_x, Inches(1.8), rec1_w, Inches(0.3),
            text="Addresses: Low Perceived Returns to AI Innovation",
            font_size=10, font_color=RED_ACCENT, bold=True, italic=True, font_name="Calibri")

# Description
r1_tb = add_rich_textbox(slide4, rec1_x, Inches(2.15), rec1_w, Inches(2.0))
tf_r1 = r1_tb.text_frame
add_paragraph(tf_r1, "The Government of India should pre-commit to procuring AI solutions in a "
              "limited set of high-priority public-sector domains. An advance market commitment (AMC) "
              "would credibly signal downstream demand, reducing uncertainty around market access "
              "for compute-intensive innovators.",
              font_size=10.5, font_color=DARK_GRAY, space_after=6)

add_paragraph(tf_r1, "Priority Use Cases:", font_size=10.5, font_color=DARK_BLUE, bold=True, space_after=3)
add_bullet_paragraph(tf_r1, "OCR of legacy government records", font_size=10, font_color=DARK_GRAY, space_after=2)
add_bullet_paragraph(tf_r1, "Multilingual citizen-facing interfaces", font_size=10, font_color=DARK_GRAY, space_after=2)
add_bullet_paragraph(tf_r1, "Forecasting for statistics, weather, and public service delivery",
                     font_size=10, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_r1, "Precedent: A $1.5B AMC for pneumococcal vaccines led to development of "
              "3 vaccines for LMICs, demonstrating the power of demand-side commitments.",
              font_size=9.5, font_color=MED_GRAY, italic=True, space_after=2)

# Implementation steps for Rec 1
impl1_bg = add_rounded_shape(slide4, rec1_x, Inches(4.3), rec1_w, Inches(2.4),
                              fill_color=VERY_LIGHT, line_color=MID_BLUE)
add_textbox(slide4, rec1_x + Inches(0.15), Inches(4.35), Inches(3), Inches(0.35),
            text="Implementation Steps", font_size=11, font_color=MID_BLUE, bold=True)

impl1_tb = add_rich_textbox(slide4, rec1_x + Inches(0.15), Inches(4.7), rec1_w - Inches(0.3), Inches(1.9))
tf_i1 = impl1_tb.text_frame

add_paragraph(tf_i1, "Step 1: Identify Priority Use Cases", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i1, "Select public-sector problem areas with clear AI applicability and "
              "measurable outputs.", font_size=9.5, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_i1, "Step 2: Define AMC Terms & Performance Criteria", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i1, "For each use case, issue an AMC specifying minimum performance "
              "standards and procurement conditions.", font_size=9.5, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_i1, "Step 3: Award Contracts & Monitor Outcomes", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i1, "Projects meeting thresholds receive longer-term procurement contracts. "
              "Ongoing monitoring tracks utilization patterns and need for adaptation.",
              font_size=9.5, font_color=DARK_GRAY, space_after=2)


# ── Recommendation 2 (right half) ──
rec2_x = Inches(6.85)
rec2_w = Inches(6.1)

# Header box
r2_header = add_rounded_shape(slide4, rec2_x, Inches(1.15), rec2_w, Inches(0.55),
                               fill_color=MID_BLUE)
add_textbox(slide4, rec2_x + Inches(0.15), Inches(1.2), rec2_w - Inches(0.3), Inches(0.45),
            text="2  Time-Bound GPU Provisioning Guarantees",
            font_size=14, font_color=WHITE, bold=True, font_name="Calibri")

# Addresses label
add_textbox(slide4, rec2_x, Inches(1.8), rec2_w, Inches(0.3),
            text="Addresses: Administrative Frictions in Access",
            font_size=10, font_color=RED_ACCENT, bold=True, italic=True, font_name="Calibri")

# Description
r2_tb = add_rich_textbox(slide4, rec2_x, Inches(2.15), rec2_w, Inches(2.0))
tf_r2 = r2_tb.text_frame
add_paragraph(tf_r2, "IndiaAI Mission should introduce clear, publicly stated service-level "
              "guarantees for compute delivery after allocation approval. Provisioning within a fixed "
              "window, with automatic escalation or fallback if timelines are breached, would directly "
              "address the credibility problem depressing demand.",
              font_size=10.5, font_color=DARK_GRAY, space_after=6)

add_paragraph(tf_r2, "Expected Impact:", font_size=10.5, font_color=DARK_BLUE, bold=True, space_after=3)
add_bullet_paragraph(tf_r2, "Reduces execution risk for firms operating under tight development cycles",
                     font_size=10, font_color=DARK_GRAY, space_after=2)
add_bullet_paragraph(tf_r2, "Allows alignment of compute access with project timelines",
                     font_size=10, font_color=DARK_GRAY, space_after=2)
add_bullet_paragraph(tf_r2, "Raises expected returns to AI experimentation through CCF",
                     font_size=10, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_r2, "Interaction: Works in tandem with Rec. 1 \u2014 demand-side pull (AMC) is "
              "only effective if supply-side delivery is credible.",
              font_size=9.5, font_color=MED_GRAY, italic=True, space_after=2)

# Implementation steps for Rec 2
impl2_bg = add_rounded_shape(slide4, rec2_x, Inches(4.3), rec2_w, Inches(2.4),
                              fill_color=VERY_LIGHT, line_color=MID_BLUE)
add_textbox(slide4, rec2_x + Inches(0.15), Inches(4.35), Inches(3), Inches(0.35),
            text="Implementation Steps", font_size=11, font_color=MID_BLUE, bold=True)

impl2_tb = add_rich_textbox(slide4, rec2_x + Inches(0.15), Inches(4.7), rec2_w - Inches(0.3), Inches(1.9))
tf_i2 = impl2_tb.text_frame

add_paragraph(tf_i2, "Step 1: Establish Service-Level Commitments", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i2, "Publicly commit to a fixed provisioning window (e.g., 2\u20134 weeks) "
              "for compute access following allocation approval.",
              font_size=9.5, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_i2, "Step 2: Create Escalation Mechanisms", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i2, "Automatic escalation protocols for breached timelines, including "
              "fallback allocation from alternative empanelled providers.",
              font_size=9.5, font_color=DARK_GRAY, space_after=5)

add_paragraph(tf_i2, "Step 3: Publish Performance Metrics", font_size=10,
              font_color=DARK_BLUE, bold=True, space_after=1)
add_paragraph(tf_i2, "Transparently report provisioning timelines to build public credibility "
              "and incentivize continuous improvement.",
              font_size=9.5, font_color=DARK_GRAY, space_after=2)

# ── Evaluation criteria bar at bottom ──
eval_y = Inches(6.85)
eval_bg = add_shape(slide4, Inches(0.5), eval_y, Inches(12.5), Inches(0.45),
                     fill_color=LIGHT_BLUE, line_color=MID_BLUE)
add_textbox(slide4, Inches(0.65), eval_y + Inches(0.03), Inches(12.2), Inches(0.38),
            text="Both recommendations rank favorably across three evaluative dimensions: "
                 "Technical Correctness  |  Administrative Feasibility  |  Political Supportability",
            font_size=10.5, font_color=DARK_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════

output_path = "/home/user/sypa/Policy Brief - Pillai Bharadwaj.pptx"
prs.save(output_path)
print(f"Policy Brief saved to: {output_path}")
